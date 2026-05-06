"""
==============================================================
GENERADOR DE PRESENTACIONES EJECUTIVAS v2.0 — RV4 Lotificaciones
Colores RV4: Navy #053D57 | Amarillo #FFCA08 | Naranja #F8931D
Incluye: cartera real, saldo inicial correcto, gráficas,
         desistimientos, alertas, análisis CFO local + OpenAI
==============================================================
"""
import io, os, logging
from datetime import datetime, date
from typing import Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)

# ── Paleta ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x05,0x3D,0x57); YELLOW = RGBColor(0xFF,0xCA,0x08)
ORANGE = RGBColor(0xF8,0x93,0x1D); WHITE  = RGBColor(0xFF,0xFF,0xFF)
GRAY   = RGBColor(0x33,0x33,0x33); LGRAY  = RGBColor(0xF2,0xF2,0xF2)
GREEN  = RGBColor(0x1A,0x85,0x38); RED    = RGBColor(0xC0,0x39,0x2B)
TEAL   = RGBColor(0x1A,0x75,0x9F); DARK_R = RGBColor(0x7B,0x24,0x1C)

SW = Inches(13.33); SH = Inches(7.5)

MESES_ES = {1:"Enero",2:"Febrero",3:"Marzo",4:"Abril",5:"Mayo",6:"Junio",
            7:"Julio",8:"Agosto",9:"Septiembre",10:"Octubre",11:"Noviembre",12:"Diciembre"}
MESES_C  = {1:"Ene",2:"Feb",3:"Mar",4:"Abr",5:"May",6:"Jun",
            7:"Jul",8:"Ago",9:"Sep",10:"Oct",11:"Nov",12:"Dic"}

# ── Mapeo proyectos (empresa_cartera = nombre en ov_cartera) ──────────────────
PROYECTOS_RV4 = [
    {"nombre":"Hacienda Jumay",               "sociedad":"EFICIENCIA URBANA","empresa_sap":"SBO_EFICIENCIA_URBANA","ec":"Eficiencia Urbana"},
    {"nombre":"La Ceiba",                     "sociedad":"SER GEN CCC",      "empresa_sap":"SBO_SER_GEN_CCC",     "ec":"Servicios Generales"},
    {"nombre":"Hacienda el Sol",              "sociedad":"ROSSIO",           "empresa_sap":"SBO_ROSSIO",          "ec":"Rossio"},
    {"nombre":"Oasis Zacapa",                 "sociedad":"FRUGALEX",         "empresa_sap":"SBO_FRUGALEX",        "ec":"Frugalex"},
    {"nombre":"Cañadas de Jalapa",            "sociedad":"OTTAVIA",          "empresa_sap":"SBO_OTTAVIA",         "ec":"Ottavia"},
    {"nombre":"Condado Jutiapa",              "sociedad":"UTILICA",          "empresa_sap":"SBO_UTILICA",         "ec":"Utilica"},
    {"nombre":"Club Campestre Jumay",         "sociedad":"TEZZOLI",          "empresa_sap":"SBO_TEZZOLI",         "ec":"Tezzoli"},
    {"nombre":"Club del Bosque",              "sociedad":"URBIVA",           "empresa_sap":"SBO_URBIVA_2",        "ec":"Urbiva 2"},
    {"nombre":"Club Residencial El Progreso", "sociedad":"GARBATELLA",       "empresa_sap":"SBO_GARBATELLA",      "ec":"Garbatella"},
    {"nombre":"Arboleada Santa Elena",        "sociedad":"CAPIPOS",          "empresa_sap":"SBO_CAPIPOS",         "ec":"Capipos"},
    {"nombre":"Hacienda Santa Lucia",         "sociedad":"OVEST",            "empresa_sap":"SBO_OVEST",           "ec":None},
    {"nombre":"Hacienda El Cafetal Fase I",   "sociedad":"CORCOLLE",         "empresa_sap":"SBO_CORCOLLE",        "ec":"Corcolle"},
    {"nombre":"Hacienda El Cafetal Fase II",  "sociedad":"LEOFRENI",         "empresa_sap":"SBO_LEOFRENI",        "ec":None},
    {"nombre":"Hacienda El Cafetal Fase III", "sociedad":"GIBRALEON",        "empresa_sap":"SBO_GIBRALEON",       "ec":"Gibraleon"},
    {"nombre":"Hacienda El Cafetal Fase IV",  "sociedad":"TALOCCI",          "empresa_sap":"SBO_TALOCCI",         "ec":None},
    {"nombre":"Celajes De Tecpan",            "sociedad":"VILET",            "empresa_sap":"SBO_VILET",           "ec":None},
]

# ── Formato ───────────────────────────────────────────────────────────────────
def _q(v):
    try:
        v=float(v or 0)
        if abs(v)>=1_000_000: return f"Q {v/1_000_000:,.2f}M"
        if abs(v)>=1_000:     return f"Q {v/1_000:,.1f}K"
        return f"Q {v:,.0f}"
    except: return "Q 0"

def _n(v):
    try: return f"{int(v or 0):,}"
    except: return "0"

def _pct(a,b):
    try: return round(float(a or 0)/float(b or 1)*100,1)
    except: return 0.0

def _ud(m,y):
    import calendar
    return f"{calendar.monthrange(y,m)[1]} de {MESES_ES[m]} de {y}"

# ── Dibujo ────────────────────────────────────────────────────────────────────
def _bg(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def _box(sl,x,y,w,h,fill,line=None):
    s=sl.shapes.add_shape(1,x,y,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line: s.line.color.rgb=line
    else: s.line.fill.background()
    return s

def _t(sl,txt,x,y,w,h,sz=11,bold=False,col=WHITE,al=PP_ALIGN.LEFT,ital=False):
    tb=sl.shapes.add_textbox(x,y,w,h)
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=al
    r=p.add_run(); r.text=str(txt or "")
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=ital
    r.font.color.rgb=col; r.font.name="Calibri"
    return tb

def _hdr(sl,title,sub="",bg=NAVY):
    _box(sl,0,0,SW,Inches(1.05),bg)
    _t(sl,title,Inches(0.35),Inches(0.07),SW-Inches(0.7),Inches(0.58),sz=20,bold=True,col=WHITE)
    if sub: _t(sl,sub,Inches(0.35),Inches(0.62),SW-Inches(0.7),Inches(0.38),sz=11,col=YELLOW)

def _ftr(sl,txt="RV4 Lotificaciones — Confidencial"):
    _box(sl,0,SH-Inches(0.27),SW,Inches(0.27),NAVY)
    _t(sl,txt,Inches(0.3),SH-Inches(0.25),SW-Inches(0.6),Inches(0.23),sz=8,col=LGRAY)

def _kpi(sl,x,y,w,h,lbl,val,bg=NAVY,vc=YELLOW,sub=""):
    _box(sl,x,y,w,h,bg)
    _t(sl,val,x,y+Inches(0.07),w,Inches(0.58),sz=21,bold=True,col=vc,al=PP_ALIGN.CENTER)
    _t(sl,lbl,x,y+Inches(0.62),w,Inches(0.28),sz=9,col=WHITE,al=PP_ALIGN.CENTER)
    if sub: _t(sl,sub,x,y+Inches(0.85),w,Inches(0.2),sz=8,col=LGRAY,al=PP_ALIGN.CENTER)

def _tbl(sl,hdrs,rows,x,y,w,h,cw_p=None):
    if not rows: return
    nc=len(hdrs)
    cw=[int(w*p) for p in cw_p] if cw_p else [w//nc]*nc
    rh=min(Inches(0.34),h//(len(rows)+1))
    for ci,hd in enumerate(hdrs):
        cx=x+sum(cw[:ci])
        _box(sl,cx,y,cw[ci],rh,NAVY)
        _t(sl,hd,cx+Inches(0.03),y+Inches(0.03),cw[ci]-Inches(0.06),rh-Inches(0.04),
           sz=9,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    for ri,row in enumerate(rows):
        bg=LGRAY if ri%2==0 else WHITE
        for ci,cell in enumerate(row):
            cx=x+sum(cw[:ci]); ry=y+(ri+1)*rh
            _box(sl,cx,ry,cw[ci],rh,bg)
            is_q="Q" in str(cell or "") or isinstance(cell,(int,float))
            _t(sl,str(cell or ""),cx+Inches(0.03),ry+Inches(0.02),
               cw[ci]-Inches(0.06),rh-Inches(0.04),sz=9,col=GRAY,
               al=PP_ALIGN.RIGHT if is_q else PP_ALIGN.LEFT)

def _bar(sl,x,y,w,h,cats,series,title=""):
    cd=ChartData(); cd.categories=cats
    for nm,vals in series: cd.add_series(nm,vals)
    ch=sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,w,h,cd).chart
    ch.has_legend=len(series)>1; ch.has_title=bool(title)
    if title:
        ch.chart_title.text_frame.text=title
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size=Pt(9)
    pal=[NAVY,GREEN,RED,ORANGE,TEAL]
    for i,s in enumerate(ch.series):
        s.format.fill.solid(); s.format.fill.fore_color.rgb=pal[i%len(pal)]
    return ch

# ── Consultas BD (índices posicionales para evitar conflictos Row) ────────────
def _si(db, soc, m, y):
    from sqlalchemy import text as T
    pi = db.execute(T('SELECT COALESCE(SUM(monto),0), MAX(anio*100+mes) FROM flujos_saldo_inicial WHERE sociedad=:s'), {'s': soc}).fetchone()
    pv = float(pi[0] or 0); mp = pi[1]
    if not mp:
        r = db.execute(T('SELECT COALESCE(SUM(monto_ingreso-monto_egreso),0) FROM flujos_efectivo WHERE sociedad=:s AND (anio<:y OR (anio=:y AND mes<:m))'), {'s':soc,'y':y,'m':m}).fetchone()
        return float(r[0] or 0)
    pa = int(mp)//100; pm2 = int(mp)%100
    ia = pa if pm2<12 else pa+1; im = pm2+1 if pm2<12 else 1
    r = db.execute(T('SELECT COALESCE(SUM(monto_ingreso-monto_egreso),0) FROM flujos_efectivo WHERE sociedad=:s AND (anio>:ia OR (anio=:ia AND mes>=:im)) AND (anio<:y OR (anio=:y AND mes<:m))'), {'s':soc,'ia':ia,'im':im,'y':y,'m':m}).fetchone()
    return pv + float(r[0] or 0)

def _flujo(db, soc, m, y):
    from sqlalchemy import text as T
    r = db.execute(T('SELECT COALESCE(SUM(monto_ingreso),0), COALESCE(SUM(monto_egreso),0) FROM flujos_efectivo WHERE sociedad=:s AND anio=:y AND mes=:m'), {'s':soc,'y':y,'m':m}).fetchone()
    i = float(r[0] or 0); e = float(r[1] or 0); si = _si(db,soc,m,y)
    di = db.execute(T("SELECT COALESCE(nombre_categoria,seccion,tipo_transaccion,'Ingreso') AS concepto, COALESCE(SUM(monto_ingreso),0) AS monto FROM flujos_efectivo WHERE sociedad=:s AND anio=:y AND mes=:m AND monto_ingreso>0 GROUP BY concepto ORDER BY monto DESC LIMIT 8"), {'s':soc,'y':y,'m':m}).fetchall()
    de = db.execute(T("SELECT COALESCE(nombre_categoria,seccion,tipo_transaccion,'Egreso') AS concepto, COALESCE(SUM(monto_egreso),0) AS monto FROM flujos_efectivo WHERE sociedad=:s AND anio=:y AND mes=:m AND monto_egreso>0 GROUP BY concepto ORDER BY monto DESC LIMIT 8"), {'s':soc,'y':y,'m':m}).fetchall()
    return {'i':i,'e':e,'si':si,'sf':si+i-e,'di':[{'c':x[0],'v':float(x[1])} for x in di],'de':[{'c':x[0],'v':float(x[1])} for x in de]}

def _flujo_acum(db, soc, m1, m2, y):
    from sqlalchemy import text as T
    rows = db.execute(T('SELECT mes, COALESCE(SUM(monto_ingreso),0), COALESCE(SUM(monto_egreso),0) FROM flujos_efectivo WHERE sociedad=:s AND anio=:y AND mes BETWEEN :m1 AND :m2 GROUP BY mes ORDER BY mes'), {'s':soc,'y':y,'m1':m1,'m2':m2}).fetchall()
    ti = sum(float(r[1]) for r in rows); te = sum(float(r[2]) for r in rows)
    si = _si(db,soc,m1,y)
    return {'i':ti,'e':te,'si':si,'sf':si+ti-te}

def _inv(db, esp):
    from sqlalchemy import text as T
    r = db.execute(T("SELECT COUNT(*), COUNT(*) FILTER(WHERE estatus='DISPONIBLE'), COUNT(*) FILTER(WHERE estatus IN('VENTA','RESERVADO')), COUNT(*) FILTER(WHERE estatus='BLOQUEADO'), COALESCE(SUM(precio_final) FILTER(WHERE estatus IN('VENTA','RESERVADO')),0) FROM lotes l JOIN proyectos p ON p.id=l.proyecto_id WHERE p.empresa_sap=:e"), {'e':esp}).fetchone()
    total=int(r[0] or 0); disp=int(r[1] or 0); vend=int(r[2] or 0)
    mz = db.execute(T("SELECT l.manzana, COUNT(*), COUNT(*) FILTER(WHERE l.estatus='DISPONIBLE'), COUNT(*) FILTER(WHERE l.estatus IN('VENTA','RESERVADO')) FROM lotes l JOIN proyectos p ON p.id=l.proyecto_id WHERE p.empresa_sap=:e AND l.manzana IS NOT NULL GROUP BY l.manzana ORDER BY l.manzana LIMIT 12"), {'e':esp}).fetchall()
    return {'t':total,'d':disp,'v':vend,'b':int(r[3] or 0),'val':float(r[4] or 0),'abs':_pct(vend,total),'mz':[{'mz':x[0],'t':int(x[1]),'d':int(x[2]),'v':int(x[3]),'a':_pct(int(x[3]),int(x[1]))} for x in mz]}

def _cart(db, ec):
    from sqlalchemy import text as T
    empty = {'t':0,'mora':0,'tm':0,'cl':0,'c30':0,'c60':0,'c90':0,'det':[]}
    if not ec: return empty
    sql_cart = (
        'SELECT '
        "COALESCE(SUM(CASE WHEN line_status='O' THEN saldo_pendiente ELSE 0 END),0), "
        "COALESCE(SUM(CASE WHEN line_status='O' AND fecha_programada_cobro<CURRENT_DATE THEN saldo_pendiente ELSE 0 END),0), "
        "COUNT(DISTINCT CASE WHEN line_status='O' THEN card_code END), "
        "COALESCE(SUM(CASE WHEN line_status='O' AND fecha_programada_cobro BETWEEN CURRENT_DATE AND CURRENT_DATE+30 THEN saldo_pendiente ELSE 0 END),0), "
        "COALESCE(SUM(CASE WHEN line_status='O' AND fecha_programada_cobro BETWEEN CURRENT_DATE AND CURRENT_DATE+60 THEN saldo_pendiente ELSE 0 END),0), "
        "COALESCE(SUM(CASE WHEN line_status='O' AND fecha_programada_cobro BETWEEN CURRENT_DATE AND CURRENT_DATE+90 THEN saldo_pendiente ELSE 0 END),0) "
        "FROM ov_cartera WHERE tipo_linea IN('BB','S') AND empresa=:e"
    )
    r = db.execute(T(sql_cart), {'e':ec}).fetchone()
    t=float(r[0] or 0); mora=float(r[1] or 0)
    sql_det = (
        'SELECT oc.card_name, COALESCE(oc.referencia_manzana_lote,\'\') AS lote, '
        "SUM(CASE WHEN oc.line_status='O' AND oc.fecha_programada_cobro<CURRENT_DATE THEN oc.saldo_pendiente ELSE 0 END) AS mv, "
        'MAX(oc.slp_name) AS asesor '
        "FROM ov_cartera oc WHERE oc.empresa=:e AND oc.tipo_linea IN('BB','S') "
        "AND oc.line_status='O' AND oc.fecha_programada_cobro<CURRENT_DATE "
        'GROUP BY oc.card_code,oc.card_name,oc.referencia_manzana_lote '
        "HAVING SUM(CASE WHEN oc.line_status='O' AND oc.fecha_programada_cobro<CURRENT_DATE THEN oc.saldo_pendiente ELSE 0 END)>0 "
        'ORDER BY mv DESC LIMIT 8'
    )
    det = db.execute(T(sql_det), {'e':ec}).fetchall()
    return {'t':t,'mora':mora,'tm':_pct(mora,t),'cl':int(r[2] or 0),'c30':float(r[3] or 0),'c60':float(r[4] or 0),'c90':float(r[5] or 0),'det':[{'n':x[0][:28],'l':x[1][:15],'v':float(x[2]),'a':(x[3] or '')[:18]} for x in det]}

def _vhist(db, esp):
    from sqlalchemy import text as T
    rows = db.execute(T("SELECT EXTRACT(MONTH FROM l.fecha_venta)::INT, COUNT(*), COALESCE(SUM(l.precio_final),0) FROM lotes l JOIN proyectos p ON p.id=l.proyecto_id WHERE p.empresa_sap=:e AND l.estatus IN('VENTA','RESERVADO') AND l.fecha_venta IS NOT NULL AND l.fecha_venta>=(CURRENT_DATE-INTERVAL '12 months') GROUP BY EXTRACT(MONTH FROM l.fecha_venta)::INT ORDER BY 1"), {'e':esp}).fetchall()
    return [type('R',(),{'m':int(x[0]),'cnt':int(x[1]),'val':float(x[2])})() for x in rows]

def _desist(db, esp):
    from sqlalchemy import text as T
    try:
        r = db.execute(T('SELECT COUNT(*), COALESCE(SUM(pagado_capital),0), COALESCE(SUM(precio_final),0) FROM desistimientos d JOIN proyectos p ON p.id=d.proyecto_id WHERE p.empresa_sap=:e'), {'e':esp}).fetchone()
        return {'t':int(r[0] or 0),'mp':float(r[1] or 0),'mo':float(r[2] or 0),'h':[]}
    except:
        try: db.rollback()
        except: pass
        return {'t':0,'mp':0,'mo':0,'h':[]}

def _vmes(db, esp, m, y):
    from sqlalchemy import text as T
    r = db.execute(T("SELECT COUNT(*), COALESCE(SUM(l.precio_final),0) FROM lotes l JOIN proyectos p ON p.id=l.proyecto_id WHERE p.empresa_sap=:e AND l.estatus IN('VENTA','RESERVADO') AND EXTRACT(YEAR FROM l.fecha_venta)=:y AND EXTRACT(MONTH FROM l.fecha_venta)=:m AND l.fecha_venta IS NOT NULL"), {'e':esp,'y':y,'m':m}).fetchone()
    return {'c':int(r[0] or 0),'v':float(r[1] or 0)}

# ── Análisis ──────────────────────────────────────────────────────────────────
def _alertas(inv,fl,cart):
    al=[]
    if cart["tm"]>15: al.append(f"Mora elevada {cart['tm']:.1f}% — gestión de cobro urgente")
    if inv["abs"]<30 and inv["t"]>50: al.append(f"Absorción baja {inv['abs']}% — revisar estrategia comercial")
    if fl["sf"]<0: al.append("Saldo final negativo — evaluar flujo próximo mes")
    if fl["e"]>fl["i"]*1.2 and fl["i"]>0: al.append("Egresos superan ingresos >20%")
    return al

def _cfo(tipo,d,nombre,m,y):
    me=MESES_ES[m]
    if tipo=="flujo":
        i=d["i"]; e=d["e"]; n=i-e; sf=d["sf"]
        t=f"{nombre} cerró {me} {y} con ingresos {_q(i)} y egresos {_q(e)}, "
        t+=(f"flujo neto positivo de {_q(n)}. Saldo final {_q(sf)} refleja posición de liquidez " + ("sólida." if sf>500000 else "adecuada.")) if n>=0 else f"resultado neto negativo de {_q(abs(n))}. Se recomienda revisar el calendario de pagos del próximo mes."
        return t
    if tipo=="inventario":
        ab=d["abs"]; v=d["v"]; t_=d["t"]
        t=f"{nombre}: {v} de {t_} lotes vendidos ({ab:.1f}% absorción), valor acumulado {_q(d['val'])}. "
        t+=("Excelente avance de ventas, proyectando cierre en corto plazo." if ab>=70 else "Ritmo moderado, reforzar acciones comerciales." if ab>=40 else "Absorción por debajo del objetivo, revisar estrategia de precios y promoción.")
        return t
    if tipo=="cartera":
        t=d["t"]; mora=d["mora"]; tm=d["tm"]; cl=d["cl"]
        s=f"Cartera de {nombre}: {_q(t)}, {cl} clientes activos. "
        s+=("Tasa de mora saludable." if tm<5 else f"Tasa de mora {tm:.1f}% en nivel de alerta moderada." if tm<15 else f"⚠ Tasa de mora {tm:.1f}% requiere atención inmediata.")
        s+=f" Proyección cobros 30d: {_q(d['c30'])}."
        return s
    if tipo=="consol":
        s=f"Portafolio de {d['np']} proyectos: ingresos {_q(d['i'])}, egresos {_q(d['e'])}. "
        s+=(f"Flujo neto positivo {_q(d['n'])}." if d['n']>=0 else f"Flujo neto negativo {_q(abs(d['n']))} — revisar proyectos de menor desempeño.")
        s+=f" Cartera total {_q(d['cart'])}, mora ponderada {d['tm']:.1f}%."
        return s
    return ""

def _oai(prompt,key,modelo="gpt-4o-mini"):
    if not key: return ""
    try:
        from openai import OpenAI
        r=OpenAI(api_key=key).chat.completions.create(
            model=modelo,
            messages=[{"role":"system","content":"Analista financiero ejecutivo, empresa inmobiliaria guatemalteca. Español, 4 oraciones máx, sin bullets."},
                      {"role":"user","content":prompt}],
            max_tokens=300,temperature=0.3)
        return r.choices[0].message.content.strip()
    except Exception as ex:
        logger.warning(f"[OpenAI] {ex}"); return ""

def _ana(tipo,d,nombre,m,y,key=""):
    if key:
        prompts={"flujo":f"Flujo de {nombre} {MESES_ES[m]} {y}: Ingresos {_q(d.get('i',0))}, Egresos {_q(d.get('e',0))}, Saldo final {_q(d.get('sf',0))}.",
                 "inventario":f"Ventas de {nombre}: {d.get('v',0)}/{d.get('t',1)} lotes, {d.get('abs',0):.1f}% absorción, valor {_q(d.get('val',0))}.",
                 "cartera":f"Cartera {nombre}: {_q(d.get('t',0))}, mora {d.get('tm',0):.1f}%, {d.get('cl',0)} clientes.",
                 "consol":f"Portafolio {d.get('np',0)} proyectos {MESES_ES[m]} {y}: Ing {_q(d.get('i',0))}, Egr {_q(d.get('e',0))}, Cartera {_q(d.get('cart',0))}, Mora {d.get('tm',0):.1f}%."}
        result=_oai(prompts.get(tipo,""),key)
        if result: return result
    return _cfo(tipo,d,nombre,m,y)

# ── SLIDES ────────────────────────────────────────────────────────────────────
def _sl_portada(prs,titulo,sub,periodo):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl,NAVY)
    _box(sl,0,SH-Inches(0.62),SW,Inches(0.62),YELLOW)
    _box(sl,Inches(0.45),Inches(0.38),Inches(2.1),Inches(0.55),WHITE)
    _t(sl,"IRV4",Inches(0.45),Inches(0.38),Inches(2.1),Inches(0.55),sz=24,bold=True,col=NAVY,al=PP_ALIGN.CENTER)
    _t(sl,titulo,Inches(1.0),Inches(2.0),Inches(11.3),Inches(1.0),sz=36,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    _t(sl,sub,Inches(1.0),Inches(3.1),Inches(11.3),Inches(0.7),sz=26,col=YELLOW,al=PP_ALIGN.CENTER)
    _t(sl,periodo,Inches(0.4),SH-Inches(0.58),Inches(12.5),Inches(0.54),sz=14,bold=True,col=NAVY,al=PP_ALIGN.CENTER)

def _sl_sep(prs,nombre,sub=""):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl,NAVY); _box(sl,0,Inches(3.05),SW,Inches(0.05),YELLOW)
    _t(sl,nombre.upper(),Inches(0.8),Inches(2.2),Inches(11.7),Inches(0.9),sz=40,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    if sub: _t(sl,sub,Inches(0.8),Inches(3.2),Inches(11.7),Inches(0.5),sz=15,col=YELLOW,al=PP_ALIGN.CENTER)
    _ftr(sl)

def _sl_flujo(prs,nombre,m,y,d,ana="",acum=False):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    per=f"ENERO — {MESES_ES[m].upper()} {y}" if acum else f"{MESES_ES[m].upper()} {y}"
    _hdr(sl,f"FLUJO DE EFECTIVO — {nombre.upper()}",f"{per} | Cifras en quetzales")
    kw=Inches(2.92); kh=Inches(1.08); ky=Inches(1.1)
    for i,(lb,vl,col) in enumerate([("SALDO INICIAL",_q(d["si"]),TEAL),("INGRESOS",_q(d["i"]),GREEN),("EGRESOS",_q(d["e"]),RED),("SALDO FINAL",_q(d["sf"]),ORANGE if d["sf"]>=0 else RED)]):
        _kpi(sl,Inches(0.22)+i*(kw+Inches(0.13)),ky,kw,kh,lb,vl,col)
    # Tablas
    _box(sl,Inches(0.22),Inches(2.32),Inches(5.9),Inches(0.3),NAVY)
    _t(sl,"INGRESOS",Inches(0.22),Inches(2.32),Inches(5.9),Inches(0.3),sz=10,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    _tbl(sl,["Concepto","Monto"],[(x["c"][:28],_q(x["v"])) for x in d.get("di",[])],Inches(0.22),Inches(2.62),Inches(5.9),Inches(2.5),[0.72,0.28])
    _box(sl,Inches(7.22),Inches(2.32),Inches(5.88),Inches(0.3),NAVY)
    _t(sl,"EGRESOS",Inches(7.22),Inches(2.32),Inches(5.88),Inches(0.3),sz=10,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    _tbl(sl,["Concepto","Monto"],[(x["c"][:28],_q(x["v"])) for x in d.get("de",[])],Inches(7.22),Inches(2.62),Inches(5.88),Inches(2.5),[0.72,0.28])
    if ana:
        _box(sl,Inches(0.22),Inches(5.25),SW-Inches(0.44),Inches(1.45),LGRAY)
        _t(sl,"📊 ANÁLISIS EJECUTIVO",Inches(0.32),Inches(5.3),Inches(5),Inches(0.28),sz=9,bold=True,col=NAVY)
        _t(sl,ana,Inches(0.32),Inches(5.58),SW-Inches(0.7),Inches(1.0),sz=10,col=GRAY)
    _ftr(sl,f"RV4 | {nombre} | {per}")

def _sl_inv(prs,nombre,m,y,inv,vh,ds,ana=""):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    _hdr(sl,f"AVANCE FINANCIERO — {nombre.upper()}",f"Al {_ud(m,y)} | Cifras en quetzales")
    kw=Inches(2.42); kh=Inches(1.08); ky=Inches(1.1)
    for i,(lb,vl,col) in enumerate([("TOTAL LOTES",_n(inv["t"]),NAVY),("DISPONIBLES",_n(inv["d"]),GREEN),("VENDIDOS",_n(inv["v"]),ORANGE),(f"ABSORCIÓN",f"{inv['abs']:.1f}%",TEAL),("DESISTIMIENTOS",_n(ds["t"]),RED)]):
        _kpi(sl,Inches(0.17)+i*(kw+Inches(0.12)),ky,kw,kh,lb,vl,col)
    _box(sl,Inches(0.17),Inches(2.3),SW-Inches(0.34),Inches(0.3),YELLOW)
    _t(sl,f"VALOR ACUMULADO VENDIDO: {_q(inv['val'])}",Inches(0.3),Inches(2.3),SW-Inches(0.6),Inches(0.3),sz=12,bold=True,col=NAVY,al=PP_ALIGN.CENTER)
    if inv["mz"]:
        _tbl(sl,["Manzana","Total","Disp.","Vendidos","Absorc."],
             [(x["mz"],x["t"],x["d"],x["v"],f"{x['a']:.1f}%") for x in inv["mz"][:8]],
             Inches(0.17),Inches(2.72),Inches(6.4),Inches(3.5),[0.3,0.18,0.18,0.18,0.16])
    if vh:
        cats=[f"{MESES_C[r.m]}" for r in vh]
        vals=[int(r.cnt) for r in vh]
        _bar(sl,Inches(6.8),Inches(2.65),Inches(6.3),Inches(3.6),cats,[("Ventas",vals)],"Ventas últimos 12 meses")
    if ds["t"]>0:
        _t(sl,f"⚠ Desistimientos históricos: {ds['t']} lotes | Capital pagado: {_q(ds['mp'])}",
           Inches(0.17),Inches(6.38),SW-Inches(0.34),Inches(0.28),sz=9,col=RED)
    if ana:
        _box(sl,Inches(0.17),Inches(6.7),SW-Inches(0.34),Inches(0.5),LGRAY)
        _t(sl,ana,Inches(0.27),Inches(6.73),SW-Inches(0.54),Inches(0.43),sz=9,col=GRAY)
    _ftr(sl,f"RV4 | {nombre} | Inventario {MESES_ES[m]} {y}")

def _sl_cart(prs,nombre,m,y,c,ana=""):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    _hdr(sl,f"ANÁLISIS DE CARTERA — {nombre.upper()}",f"Al cierre de {MESES_ES[m]} {y} | Cifras en quetzales")
    kw=Inches(2.42); kh=Inches(1.08); ky=Inches(1.1)
    mc=RED if c["tm"]>15 else (ORANGE if c["tm"]>5 else GREEN)
    for i,(lb,vl,col) in enumerate([("CARTERA TOTAL",_q(c["t"]),NAVY),("EN MORA",_q(c["mora"]),mc),("TASA MORA",f"{c['tm']:.1f}%",mc),("CLIENTES",_n(c["cl"]),TEAL),("COBRO 30D",_q(c["c30"]),GREEN)]):
        _kpi(sl,Inches(0.17)+i*(kw+Inches(0.12)),ky,kw,kh,lb,vl,col)
    _box(sl,Inches(0.17),Inches(2.3),SW-Inches(0.34),Inches(0.28),TEAL)
    _t(sl,f"PROYECCIÓN COBROS | 30d: {_q(c['c30'])}  |  60d: {_q(c['c60'])}  |  90d: {_q(c['c90'])}",
       Inches(0.3),Inches(2.3),SW-Inches(0.6),Inches(0.28),sz=11,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
    if c["det"]:
        _box(sl,Inches(0.17),Inches(2.7),SW-Inches(0.34),Inches(0.3),NAVY)
        _t(sl,"CLIENTES CON SALDO VENCIDO",Inches(0.17),Inches(2.7),SW-Inches(0.34),Inches(0.3),sz=10,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
        _tbl(sl,["Cliente","Lote","Vencido","Asesor"],
             [(x["n"],x["l"],_q(x["v"]),x["a"]) for x in c["det"]],
             Inches(0.17),Inches(3.0),SW-Inches(0.34),Inches(3.2),[0.35,0.22,0.22,0.21])
    else:
        _box(sl,Inches(0.17),Inches(2.7),SW-Inches(0.34),Inches(0.45),LGRAY)
        _t(sl,"✓  Sin clientes con saldo vencido al cierre del período",
           Inches(0.3),Inches(2.78),SW-Inches(0.6),Inches(0.32),sz=12,col=GREEN,al=PP_ALIGN.CENTER)
    if ana:
        _box(sl,Inches(0.17),Inches(6.3),SW-Inches(0.34),Inches(0.9),LGRAY)
        _t(sl,ana,Inches(0.27),Inches(6.35),SW-Inches(0.54),Inches(0.8),sz=10,col=GRAY)
    _ftr(sl,f"RV4 | {nombre} | Cartera {MESES_ES[m]} {y}")

def _sl_flujo_consol(prs,m,y,plist,acum=False):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    per=f"ENERO — {MESES_ES[m].upper()} {y}" if acum else f"{MESES_ES[m].upper()} {y}"
    _hdr(sl,"FLUJO DE EFECTIVO CONSOLIDADO",f"{per} | Cifras en quetzales")
    ti=sum(p["i"] for p in plist); te=sum(p["e"] for p in plist); tn=ti-te
    kw=Inches(2.92); kh=Inches(1.05); ky=Inches(1.1)
    for i,(lb,vl,col) in enumerate([("TOTAL INGRESOS",_q(ti),GREEN),("TOTAL EGRESOS",_q(te),RED),("RESULTADO NETO",_q(tn),NAVY if tn>=0 else RED),("SALDO FINAL",_q(sum(p["sf"] for p in plist)),ORANGE)]):
        _kpi(sl,Inches(0.22)+i*(kw+Inches(0.13)),ky,kw,kh,lb,vl,col)
    _tbl(sl,["Proyecto","Saldo Inicial","Ingresos","Egresos","Saldo Final"],
         [(p["n"][:26],_q(p["si"]),_q(p["i"]),_q(p["e"]),_q(p["sf"])) for p in plist],
         Inches(0.22),Inches(2.32),SW-Inches(0.44),Inches(4.85),[0.32,0.17,0.17,0.17,0.17])
    _ftr(sl,f"RV4 | Flujo Consolidado | {per}")

def _sl_cart_consol(prs,m,y,plist):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    _hdr(sl,"RECUPERACIÓN DE CARTERA CONSOLIDADA",f"Al cierre de {MESES_ES[m]} {y} | Cifras en quetzales")
    tc=sum(p["ct"] for p in plist); tm=sum(p["mora"] for p in plist)
    kw=Inches(3.8); kh=Inches(1.05); ky=Inches(1.1)
    mc=RED if _pct(tm,tc)>15 else (ORANGE if _pct(tm,tc)>5 else GREEN)
    for i,(lb,vl,col) in enumerate([("CARTERA TOTAL",_q(tc),NAVY),("TOTAL EN MORA",_q(tm),mc),(f"TASA MORA POND.",f"{_pct(tm,tc):.1f}%",mc)]):
        _kpi(sl,Inches(0.3)+i*(kw+Inches(0.27)),ky,kw,kh,lb,vl,col)
    _tbl(sl,["Proyecto","Cartera","Mora","Tasa Mora","Clientes"],
         [(p["n"][:26],_q(p["ct"]),_q(p["mora"]),f"{p['tm']:.1f}%",_n(p["cl"])) for p in plist],
         Inches(0.22),Inches(2.32),SW-Inches(0.44),Inches(4.85),[0.36,0.18,0.18,0.15,0.13])
    _ftr(sl,f"RV4 | Cartera Consolidada | {MESES_ES[m]} {y}")

def _sl_ventas_consol(prs,m,y,plist):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    _hdr(sl,"CONSOLIDADO DE VENTAS",f"Al cierre de {MESES_ES[m]} {y}")
    tv=sum(p["vm"] for p in plist); tval=sum(p["vv"] for p in plist)
    ttot=sum(p["t"] for p in plist); tabs=_pct(sum(p["v"] for p in plist),ttot)
    kw=Inches(2.92); kh=Inches(1.05); ky=Inches(1.1)
    for i,(lb,vl,col) in enumerate([("LOTES VENDIDOS MES",_n(tv),ORANGE),("VALOR VENTAS MES",_q(tval),NAVY),("ABSORCIÓN PORTAF.",f"{tabs:.1f}%",GREEN),("TOTAL LOTES",_n(ttot),TEAL)]):
        _kpi(sl,Inches(0.22)+i*(kw+Inches(0.13)),ky,kw,kh,lb,vl,col)
    _tbl(sl,["Proyecto","Total","Vendidos","Disp.","Absorción","Ventas Mes"],
         [(p["n"][:26],_n(p["t"]),_n(p["v"]),_n(p["d"]),f"{p['a']:.1f}%",_n(p["vm"])) for p in plist],
         Inches(0.22),Inches(2.32),SW-Inches(0.44),Inches(4.85),[0.33,0.12,0.13,0.13,0.15,0.14])
    _ftr(sl,f"RV4 | Ventas Consolidadas | {MESES_ES[m]} {y}")

def _sl_alertas(prs,m,y,adict):
    al=[(n,a) for n,als in adict.items() for a in als]
    if not al: return
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
    _hdr(sl,"⚠  ALERTAS Y PUNTOS DE ATENCIÓN",f"Período: {MESES_ES[m]} {y}",DARK_R)
    for i,(n,a) in enumerate(al[:14]):
        ry=Inches(1.18)+i*Inches(0.41)
        _box(sl,Inches(0.2),ry,SW-Inches(0.4),Inches(0.39),LGRAY if i%2==0 else WHITE)
        _t(sl,f"⚠  {n}: {a}",Inches(0.33),ry+Inches(0.05),SW-Inches(0.66),Inches(0.3),sz=10,col=DARK_R)
    _ftr(sl)

# ── Generadores principales ───────────────────────────────────────────────────
def _prs():
    p=Presentation(); p.slide_width=SW; p.slide_height=SH; return p

def generar_presentacion_proyecto(db, empresa_sap, mes, anio, openai_key=""):
    info=next((p for p in PROYECTOS_RV4 if p["empresa_sap"]==empresa_sap),None)
    if not info: raise ValueError(f"Proyecto {empresa_sap} no encontrado")
    nm=info["nombre"]; soc=info["sociedad"]; ec=info["ec"]
    fl=_flujo(db,soc,mes,anio); fa=_flujo_acum(db,soc,1,mes,anio)
    inv=_inv(db,empresa_sap); cart=_cart(db,ec)
    vh=_vhist(db,empresa_sap); ds=_desist(db,empresa_sap)
    al=_alertas(inv,fl,cart)
    prs=_prs()
    _sl_portada(prs,"Junta Directiva Lotificadoras",nm.upper(),f"Resultados a {MESES_ES[mes]} {anio}")
    _sl_sep(prs,nm)
    _sl_flujo(prs,nm,mes,anio,fl,_ana("flujo",fl,nm,mes,anio,openai_key),False)
    if mes>1: _sl_flujo(prs,nm,mes,anio,{"si":fa["si"],"i":fa["i"],"e":fa["e"],"sf":fa["sf"],"di":[],"de":[]},_ana("flujo",{"i":fa["i"],"e":fa["e"],"si":fa["si"],"sf":fa["sf"]},nm,mes,anio,openai_key),True)
    _sl_inv(prs,nm,mes,anio,inv,vh,ds,_ana("inventario",inv,nm,mes,anio,openai_key))
    _sl_cart(prs,nm,mes,anio,cart,_ana("cartera",cart,nm,mes,anio,openai_key))
    if al:
        sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,WHITE)
        _hdr(sl,f"ALERTAS — {nm.upper()}","Puntos de atención del período",DARK_R)
        for i,a in enumerate(al):
            _box(sl,Inches(0.3),Inches(1.2)+i*Inches(0.5),SW-Inches(0.6),Inches(0.45),LGRAY)
            _t(sl,f"⚠  {a}",Inches(0.45),Inches(1.25)+i*Inches(0.5),SW-Inches(0.9),Inches(0.38),sz=12,col=DARK_R)
        _ftr(sl)
    buf=io.BytesIO(); prs.save(buf); buf.seek(0); return buf.read()

def generar_presentacion_consolidada(db, mes, anio, openai_key=""):
    prs=_prs()
    _sl_portada(prs,"Junta Directiva Lotificadoras","CONSOLIDADO",f"Resultados a {MESES_ES[mes]} {anio}")
    fl_consol=[]; cart_consol=[]; inv_consol=[]; alertas_dict={}
    for p in PROYECTOS_RV4:
        fl=_flujo(db,p["sociedad"],mes,anio); fa=_flujo_acum(db,p["sociedad"],1,mes,anio)
        inv=_inv(db,p["empresa_sap"]); cart=_cart(db,p["ec"])
        vm=_vmes(db,p["empresa_sap"],mes,anio)
        al=_alertas(inv,fl,cart)
        if al: alertas_dict[p["nombre"]]=al
        fl_consol.append({"n":p["nombre"],"si":fl["si"],"i":fl["i"],"e":fl["e"],"sf":fl["sf"],"si_a":fa["si"],"i_a":fa["i"],"e_a":fa["e"],"sf_a":fa["sf"]})
        cart_consol.append({"n":p["nombre"],"ct":cart["t"],"mora":cart["mora"],"tm":cart["tm"],"cl":cart["cl"]})
        inv_consol.append({"n":p["nombre"],"t":inv["t"],"v":inv["v"],"d":inv["d"],"a":inv["abs"],"vm":vm["c"],"vv":vm["v"]})

    _sl_sep(prs,"CONSOLIDADO","Resumen de todos los proyectos")
    _sl_flujo_consol(prs,mes,anio,fl_consol,False)
    if mes>1:
        _sl_flujo_consol(prs,mes,anio,[{"n":p["n"],"si":p["si_a"],"i":p["i_a"],"e":p["e_a"],"sf":p["sf_a"]} for p in fl_consol],True)
    _sl_cart_consol(prs,mes,anio,cart_consol)
    _sl_ventas_consol(prs,mes,anio,inv_consol)
    ti=sum(p["i"] for p in fl_consol); te=sum(p["e"] for p in fl_consol)
    tc=sum(p["ct"] for p in cart_consol); tm=sum(p["mora"] for p in cart_consol)
    ana_c=_ana("consol",{"np":len(PROYECTOS_RV4),"i":ti,"e":te,"n":ti-te,"cart":tc,"tm":_pct(tm,tc)},"",mes,anio,openai_key)
    if ana_c:
        sl=prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,NAVY)
        _t(sl,"ANÁLISIS EJECUTIVO DEL PORTAFOLIO",Inches(1),Inches(1.5),Inches(11.3),Inches(0.7),sz=24,bold=True,col=WHITE,al=PP_ALIGN.CENTER)
        _box(sl,Inches(1),Inches(2.35),Inches(11.3),Inches(0.05),YELLOW)
        _t(sl,ana_c,Inches(1.2),Inches(2.6),Inches(11.0),Inches(3.5),sz=15,col=LGRAY,al=PP_ALIGN.CENTER)
        _ftr(sl)
    if alertas_dict: _sl_alertas(prs,mes,anio,alertas_dict)
    _sl_sep(prs,"DETALLE POR PROYECTO","Análisis individual")
    for p in PROYECTOS_RV4:
        fl=_flujo(db,p["sociedad"],mes,anio); fa=_flujo_acum(db,p["sociedad"],1,mes,anio)
        inv=_inv(db,p["empresa_sap"]); cart=_cart(db,p["ec"])
        vh=_vhist(db,p["empresa_sap"]); ds=_desist(db,p["empresa_sap"])
        _sl_sep(prs,p["nombre"])
        _sl_flujo(prs,p["nombre"],mes,anio,fl,_ana("flujo",fl,p["nombre"],mes,anio,openai_key),False)
        if mes>1: _sl_flujo(prs,p["nombre"],mes,anio,{"si":fa["si"],"i":fa["i"],"e":fa["e"],"sf":fa["sf"],"di":[],"de":[]},_ana("flujo",{"i":fa["i"],"e":fa["e"],"si":fa["si"],"sf":fa["sf"]},p["nombre"],mes,anio,openai_key),True)
        _sl_inv(prs,p["nombre"],mes,anio,inv,vh,ds,_ana("inventario",inv,p["nombre"],mes,anio,openai_key))
        _sl_cart(prs,p["nombre"],mes,anio,cart,_ana("cartera",cart,p["nombre"],mes,anio,openai_key))
    buf=io.BytesIO(); prs.save(buf); buf.seek(0); return buf.read()
